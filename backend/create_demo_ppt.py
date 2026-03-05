from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_demo_ppt():
    prs = Presentation()
    
    # Define Layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Set Theme Colors (Gurukul colors - Dark Theme)
    theme_bg_color = RGBColor(10, 12, 8) # #0a0c08
    text_color = RGBColor(255, 255, 255)
    accent_color = RGBColor(245, 158, 11) # Orange/Amber
    
    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Gurukul Learning Platform"
    # Format Title
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "AI-Powered Personalized Education\nAligned with NEP 2020"
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Helper to add basic content slide
    def add_slide(title_text, bullet_points):
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = accent_color
        
        body = slide.placeholders[1]
        tf = body.text_frame
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(28)
            p.font.color.rgb = text_color
            
    # 2. Problem Statement
    add_slide("The Problem", [
        "One-size-fits-all education leaves students behind",
        "Teachers are overwhelmed with administrative tasks",
        "Parents lack real-time visibility into learning progress",
        "Institutions struggle to track true learning outcomes"
    ])

    # 3. Gurukul Solution
    add_slide("The Gurukul Solution", [
        "AI-Driven Personalized Learning Paths",
        "Automated Assessments and Analytics",
        "Transparent Parent-Teacher Communication",
        "Secure, Scalable Cloud Infrastructure"
    ])

    # 4. Student Experience
    add_slide("Student Experience", [
        "Tailored content generation based on learning style",
        "Interactive AI Chatbot for 24/7 assistance",
        "Gamified progress tracking (Karma & Streaks)",
        "Daily reflection and mindfulness features"
    ])

    # 5. Teacher Experience
    add_slide("Teacher Experience", [
        "Real-time pulse of classroom comprehension",
        "Automated quiz generation and grading",
        "Early warning system for struggling students",
        "Streamlined communication channels"
    ])

    # 6. Parent Experience
    add_slide("Parent Experience", [
        "Non-intrusive dashboard for peace of mind",
        "Weekly automated progress summaries",
        "Direct visibility into test scores and teacher feedback",
        "Clear alignment with school curriculum"
    ])

    # 7. Institution Benefits
    add_slide("Institution Benefits", [
        "Unified data ecosystem across all stakeholders",
        "Data-driven strategic planning and resource allocation",
        "High security and data isolation per tenant",
        "Compliance with modern educational standards"
    ])

    # 8. NEP 2020 Alignment
    add_slide("NEP 2020 Alignment", [
        "Focus on core concepts and critical thinking over rote learning",
        "Multidisciplinary and holistic education",
        "Technology-enabled continuous assessment",
        "Respect for diversity and local context via multi-language support"
    ])
    
    # 9. Conclusion
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Thank You"
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    
    prs.save('c:/Users/pc45/Desktop/Gurukul/Gurukul_Demo_Presentation.pptx')
    print("Successfully created Gurukul_Demo_Presentation.pptx")

if __name__ == "__main__":
    create_demo_ppt()
